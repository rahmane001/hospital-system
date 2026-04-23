"""
build_presentation.py — Regenerate HMS_DApp_Presentation.pptx.

Produces a 14-slide editable deck at 1920x1080 (16:9 widescreen) matching the
HTML design at /tmp/hms-design/hms-dapp-pptx/project/HMS DApp Presentation.html.

Usage:
    python3 scripts/build_presentation.py

Optional screenshots at presentation/screenshots/shot-0[1-6]-*.png are embedded
on Slide 9. Missing files are skipped with a warning (the cell is left blank).

Dependencies:
    pip3 install --user python-pptx
"""
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --------------------------------------------------------------------------
# Layout constants — 1920x1080 at 144 DPI → 13.333" x 7.5" widescreen
# --------------------------------------------------------------------------
SLIDE_W = Inches(13.3333)
SLIDE_H = Inches(7.5)
PX = 6350  # EMU per design pixel (1 inch = 914400 EMU; 144 px = 1 inch)

REPO_ROOT = Path(__file__).resolve().parent.parent
SCREENSHOT_DIR = REPO_ROOT / "presentation" / "screenshots"
OUTPUT_PATH = REPO_ROOT / "HMS_DApp_Presentation.pptx"

def px(n): return Emu(int(n * PX))
def fpt(n): return Pt(n * 0.5)  # design px (font) → pt

# --------------------------------------------------------------------------
# Colour palette — matches the HTML CSS variables
# --------------------------------------------------------------------------
NAVY      = RGBColor(0x0d, 0x1b, 0x2a)
NAVY2     = RGBColor(0x12, 0x20, 0x33)
WHITE     = RGBColor(0xf9, 0xfa, 0xfb)
OFFWHITE  = RGBColor(0xee, 0xf0, 0xf4)
BLUE      = RGBColor(0x1a, 0x8f, 0xc2)
BLUE_LT   = RGBColor(0x6c, 0xb8, 0xda)
TEAL      = RGBColor(0x0a, 0xa6, 0x93)
TEAL_LT   = RGBColor(0x7a, 0xd2, 0xc0)
AMBER     = RGBColor(0xd9, 0xa0, 0x4a)
AMBER_DK  = RGBColor(0xbb, 0x85, 0x20)
RED       = RGBColor(0xc8, 0x4a, 0x4a)
TEXT      = RGBColor(0x0d, 0x1b, 0x2a)
MUTED     = RGBColor(0x5a, 0x6a, 0x7e)
BORDER    = RGBColor(0xd8, 0xdd, 0xe6)
ROW_ALT   = RGBColor(0xf5, 0xf7, 0xfa)

BG_BLUE_CARD  = RGBColor(0xea, 0xf5, 0xfc)
BG_TEAL_CARD  = RGBColor(0xe6, 0xf7, 0xf2)
BG_AMBER_CARD = RGBColor(0xfa, 0xf3, 0xe6)
BG_SLATE_CARD = RGBColor(0xf0, 0xf2, 0xf5)

AB_REACT_BG   = RGBColor(0x16, 0x2c, 0x44); AB_REACT_BR = RGBColor(0x24, 0x4a, 0x70)
AB_NODE_BG    = RGBColor(0x13, 0x32, 0x2a); AB_NODE_BR  = RGBColor(0x21, 0x55, 0x48)
AB_CHAIN_BG   = RGBColor(0x34, 0x28, 0x14); AB_CHAIN_BR = RGBColor(0x6b, 0x52, 0x2a)

TAG_REACT_BG  = RGBColor(0x1f, 0x3f, 0x60); TAG_REACT_FG = RGBColor(0x8e, 0xc6, 0xe0)
TAG_NODE_BG   = RGBColor(0x1a, 0x44, 0x3a); TAG_NODE_FG  = RGBColor(0x8a, 0xd6, 0xc4)
TAG_CHAIN_BG  = RGBColor(0x4a, 0x38, 0x1d); TAG_CHAIN_FG = RGBColor(0xde, 0xbb, 0x6e)

FN_COLORS = [
    BLUE, RGBColor(0x14, 0x97, 0xb5), TEAL,
    RGBColor(0x3a, 0xa8, 0x70), RGBColor(0x5a, 0xaf, 0x52), AMBER,
]
DOT_PURPLE = RGBColor(0xa0, 0x8b, 0xd4)
DOT_RED    = RGBColor(0xde, 0x8f, 0x8a)
DOT_GREEN  = RGBColor(0x8b, 0xc4, 0x8a)

CODE_BG = RGBColor(0x0a, 0x14, 0x22)
CODE_FG = RGBColor(0xd6, 0xe4, 0xf2)

# --------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# Core document metadata
_cp = prs.core_properties
_cp.author = "Esha Rahman"
_cp.last_modified_by = "Esha Rahman"
_cp.title = "HMS DApp - Hospital Management on the Blockchain"
_cp.subject = "CN6035 Mobile and Distributed Systems Coursework"
_cp.comments = ""

def add_slide(bg_color):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = bg_color
    bg.shadow.inherit = False
    return s

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if line_w is not None: sh.line.width = line_w
    if radius:
        try: sh.adjustments[0] = radius
        except Exception: pass
    return sh

def add_text(slide, x, y, w, h, text, size=18, bold=False, color=TEXT,
             font="DM Sans", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, letter_spacing=None, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        f = r.font; f.name = font
        f.size = size if isinstance(size, Pt) else Pt(size)
        f.bold = bold; f.italic = italic; f.color.rgb = color
        if letter_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(letter_spacing))
    return tb

def add_multiline(slide, x, y, w, h, runs, anchor=MSO_ANCHOR.TOP,
                  align=PP_ALIGN.LEFT, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        for rd in para_runs:
            r = p.add_run(); r.text = rd.get('text', '')
            f = r.font
            f.name = rd.get('font', 'DM Sans')
            size = rd.get('size', Pt(12))
            f.size = size if isinstance(size, Pt) else Pt(size)
            f.bold = rd.get('bold', False); f.italic = rd.get('italic', False)
            f.color.rgb = rd.get('color', TEXT)
    return tb

def set_notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(12); r.font.name = 'Calibri'

def slide_label(slide, text, color=BLUE, x=px(120), y=px(100)):
    add_text(slide, x, y, px(1000), px(34), text, size=fpt(24),
             bold=True, color=color, letter_spacing=240)

def slide_title(slide, text, color=TEXT, x=px(120), y=px(134), w=px(1680)):
    add_text(slide, x, y, w, px(90), text, size=fpt(72), bold=True,
             color=color, line_spacing=1.05)

def add_line(slide, x1, y1, x2, y2, color=BORDER, weight=Pt(2), dash=False):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = weight
    if dash:
        from pptx.oxml.ns import qn
        spPr = ln.line._get_or_add_ln()
        # set dashed style
        from lxml import etree
        prstDash = etree.SubElement(spPr, qn('a:prstDash'))
        prstDash.set('val', 'dash')
    return ln

def add_arrow(slide, x1, y1, x2, y2, color=MUTED, weight=Pt(2)):
    from lxml import etree
    from pptx.oxml.ns import qn
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = weight
    spPr = ln.line._get_or_add_ln()
    tail = etree.SubElement(spPr, qn('a:tailEnd'))
    tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    return ln

# --------------------------------------------------------------------------
# Speaker notes — expanded for 14 slides with depth of knowledge emphasis
# --------------------------------------------------------------------------
NOTES = [
    # 1. Title
    "Welcome. I'm presenting HMS DApp — a Hybrid Hospital Management Decentralised Application built for CN6035 Mobile and Distributed Systems. Over the next ten minutes I'll walk you through what the system does, how it's built, specifically how it interacts with the blockchain, and critically evaluate the design choices I made and where the system could be pushed further.",

    # 2. Overview
    "HMS DApp is a full-stack hospital platform that extends a traditional Node.js backend with a React frontend, a Solidity smart contract, and MetaMask wallet integration. The system supports four distinct user roles: Admin, Doctor, Patient, and Receptionist. The core value proposition is a tamper-proof audit layer — every booking, payment, and prescription is mirrored onto an Ethereum-compatible blockchain, so the backend alone cannot silently alter medical or financial records.",

    # 3. Architecture
    "The architecture is a three-tier hybrid. React 19 talks to a Node.js Express API via JWT-authenticated REST calls. Express owns business logic, writes to MongoDB for persistent storage, and dispatches web3.js calls to a local Ganache chain. Crucially, the React frontend can also bypass the backend and talk directly to MetaMask via window.ethereum for on-chain payments. Every mutating action therefore produces both a MongoDB write and an on-chain transaction, with the transaction hash stored back on the Mongo document so the two stores can be cross-verified.",

    # 4. Tech stack
    "On the technology choices — React 19 with Router for role-based SPA routing, recharts for the admin analytics dashboards, jsPDF for exportable bills and prescriptions. The backend is Express 4 with Mongoose schemas, JWT plus bcrypt for auth, node-cron for scheduled slot cleanup, and web3.js for contract calls. Blockchain is Solidity 0.8.19 deployed with Truffle onto Ganache at chain ID 5777. Each tier maps to a distinct architectural layer — Presentation, Application, Persistence, and Consensus — which is how I reasoned about cross-cutting concerns like auth and audit.",

    # 5. Smart contract
    "The smart contract HospitalRecords defines four structs — PatientRecord, AppointmentRecord, BillingRecord, PrescriptionRecord — each keyed by a MongoDB ObjectId string. Writes are gated by an onlyOwner modifier so only the backend's deployer account can append records; this is the centralisation I'll critique later. The exception is payBillOnChain, which is payable — a patient sends ETH directly via MetaMask, the contract marks the bill paid, forwards the ETH to the owner, and emits a PaymentReceived event. This is shown in the Solidity snippet on the right.",

    # 6. Flow (swimlane)
    "Here's the end-to-end flow across the stack. A patient POSTs to /api/appointments. Express authenticates the JWT, writes the Appointment document to MongoDB, then calls logAppointment on the contract via web3.js. Ganache mines the transaction and returns a hash — we persist that hash plus a blockchainStatus of 'logged' back onto the Mongo document. If the send fails we now persist a 'failed' status instead of swallowing the error, which I'll come back to on the evaluation slide. Finally a Notification is dispatched to the patient and the Admin Blockchain Verification page can independently call getAppointment on the contract to re-prove the record.",

    # 7. Off-chain vs on-chain payment
    "Payment is where the hybrid design really matters. The standard path — left side — is backend-mediated: the patient clicks Pay, the backend marks the bill paid in Mongo, then logBilling writes a status record on chain. The trust assumption is that the backend is honest. The on-chain path — right side — uses payBillOnChain: the patient signs a transaction in MetaMask, ETH goes from their wallet into the contract, the contract itself flips the bill status and forwards funds to the owner. No backend custody, no backend write. If the backend were compromised it still couldn't forge a payment here. That's what 'fully decentralised payment' actually means in the submission.",

    # 8. Roles table
    "This table summarises what each role can do. Admin is broadest — analytics, blockchain verification, audit log. Doctor manages their own slots and writes prescriptions. Patient books, pays, and reads. Receptionist handles beds and department allocation. Notice that MetaMask on-chain payment is patient-exclusive — because it's the patient's wallet signing. All roles see in-app notifications polled every thirty seconds, and PDF export is available to everyone who needs to archive a record externally.",

    # 9. Live screenshots
    "These are captures from the running system. Top-left is the admin overview with recharts-powered analytics for revenue, per-doctor load, diagnosis distribution, and on-chain coverage. Top-centre shows the patient booking a slot. Top-right is MetaMask mid-confirmation during an on-chain bill payment. Bottom-left is the Blockchain Verification page where the admin clicks Run Verification and cross-checks the tx hashes stored in Mongo against getAppointment and getBilling reads directly from Ganache. Bottom-centre is the audit log middleware output. Bottom-right is a jsPDF prescription export, which carries the tx hash in the footer so a paper copy is still verifiable.",

    # 10. Cross-cutting
    "Several cross-cutting capabilities tie the system together. Every mutating HTTP request passes through audit middleware that writes an AuditLog entry independently of the blockchain. Notifications fire on booking, payment, and prescription events. jsPDF exports bills and prescriptions. The MetaMask bar at the top of every page detects window.ethereum, shows the connected address and balance, and is the entry point for on-chain payments. JWT with bcrypt and a doctor-approval workflow secures every endpoint, and there's a seed script pre-wiring admin, doctor, and patient accounts to specific Ganache wallets so demos reproduce cleanly.",

    # 11. Design rationale & critical evaluation
    "Now the critical analysis — what I chose and why, where it falls short, and what I would do next. I chose Ganache over Sepolia because the development feedback loop is essentially free — no faucet, no block time — but this means my submission cannot be verified by a tutor on a public explorer; that's a trade-off. The onlyOwner modifier keeps write authority centralised on a single deployer account, which simplifies authentication but reintroduces the trust assumption the blockchain was meant to remove. I don't have an event indexer — so querying historical records at scale means iterating through mappings. Future work: migrate to Sepolia, index events via TheGraph, expand the Jest suite beyond the three coursework-critical paths, introduce a multi-sig deployer, and add input validation middleware using express-validator. These are the concrete next steps, not aspirational bullet points.",

    # 12. Baseline vs Extension
    "This slide makes the scope of my contribution explicit for the marker. The master branch is the pristine upstream baseline — a backend-only Node, Express, and MongoDB CRUD application with no blockchain, no React frontend, and no role-aware UI. The main branch is the submission: I layered a React 19 single-page app, a Solidity smart contract deployed to Ganache, web3.js integration, MetaMask payments, recharts analytics, jsPDF exports, audit middleware, notifications, and the entire Task 2 presentation bundle on top. The two branches live on the same GitHub repo so the tutor can diff them directly using the compare URL at the bottom of the slide. 112 files changed, roughly 59-thousand lines added.",

    # 13. Video
    "This slide contains the MS Stream link to the recorded demonstration. Access has been granted to all lab tutors. The recording walks through each of the four roles, an end-to-end booking with the transaction hash appearing live in the Blockchain Verification page, a MetaMask-signed payment, and the audit log.",

    # 14. References
    "In summary, HMS DApp extends a traditional CRUD application into a hybrid DApp by anchoring critical records to an immutable ledger without sacrificing usability. Here are the references used in designing and evaluating the system, in Harvard format. Thank you."
]

# --------------------------------------------------------------------------
# SLIDE 1 — Title
# --------------------------------------------------------------------------
s = add_slide(NAVY)
add_rect(s, 0, 0, px(8), SLIDE_H, fill=BLUE)
add_text(s, px(120), px(400), px(1500), px(30),
         "CN6035 — MOBILE AND DISTRIBUTED SYSTEMS",
         size=fpt(24), bold=True, color=BLUE_LT, letter_spacing=240)
add_multiline(s, px(120), px(460), px(1600), px(130), [
    [{'text': 'HMS ', 'size': fpt(96), 'bold': True, 'color': WHITE},
     {'text': 'DApp', 'size': fpt(96), 'bold': True, 'color': BLUE_LT}]
], line_spacing=1.05)
add_text(s, px(120), px(620), px(1500), px(80),
         "A Hybrid Hospital Management Decentralised Application",
         size=fpt(44), color=RGBColor(0xa6, 0xaf, 0xbb), line_spacing=1.4)
add_multiline(s, px(1100), px(900), px(720), px(100),
    [[{'text': 'Task 2 — Video Presentation', 'size': fpt(24),
       'color': RGBColor(0x66, 0x6d, 0x78)}],
     [{'text': 'Solidity · Node.js · React 19 · Ganache', 'size': fpt(24),
       'color': RGBColor(0x66, 0x6d, 0x78)}]],
    align=PP_ALIGN.RIGHT, line_spacing=1.7)
set_notes(s, NOTES[0])

# --------------------------------------------------------------------------
# SLIDE 2 — Project Overview
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "01 — OVERVIEW", BLUE)
slide_title(s, "Project Overview")

card_w = px(820); card_h = px(310); gap = px(32)
grid_top = px(310); cols = [px(120), px(120) + card_w + gap]
rows = [grid_top, grid_top + card_h + gap]

def role_card(x, y, role, items, bg, border, role_color):
    add_rect(s, x, y, card_w, card_h, fill=bg, line=border,
             line_w=Pt(1.5), radius=0.05)
    add_text(s, x + px(48), y + px(34), card_w - px(96), px(34),
             role, size=fpt(28), bold=True, color=role_color, letter_spacing=240)
    runs = []
    for it in items:
        runs.append([
            {'text': '•  ', 'size': fpt(30), 'color': role_color, 'bold': True},
            {'text': it, 'size': fpt(30), 'color': TEXT},
        ])
    add_multiline(s, x + px(48), y + px(90), card_w - px(96),
                  card_h - px(110), runs, line_spacing=1.35)

role_card(cols[0], rows[0], "ADMIN", [
    "Analytics dashboard with recharts",
    "Manage users, departments & beds",
    "Approve doctor registrations",
    "Blockchain verification page",
    "Full audit log viewer",
], BG_BLUE_CARD, BLUE_LT, BLUE)
role_card(cols[1], rows[0], "DOCTOR", [
    "Manage appointment slots",
    "View booked patients",
    "Create prescriptions (hash on-chain)",
    "View billing for own patients",
], BG_TEAL_CARD, TEAL_LT, TEAL)
role_card(cols[0], rows[1], "PATIENT", [
    "Browse doctors & book appointments",
    "View bills — pay off-chain or via MetaMask",
    "View prescriptions & notifications",
    "PDF export of bills & prescriptions",
], BG_AMBER_CARD, AMBER, AMBER_DK)
role_card(cols[1], rows[1], "RECEPTIONIST", [
    "Manage bed assignments",
    "Manage department allocations",
], BG_SLATE_CARD, RGBColor(0xc8, 0xcd, 0xd6), MUTED)
set_notes(s, NOTES[1])

# --------------------------------------------------------------------------
# SLIDE 3 — System Architecture
# --------------------------------------------------------------------------
s = add_slide(NAVY)
slide_label(s, "02 — ARCHITECTURE", BLUE_LT)
slide_title(s, "System Architecture", color=WHITE)

box_y = px(360); box_h = px(440); box_w = px(440); gap_x = px(60)
total_w = box_w * 3 + gap_x * 2
start_x = (SLIDE_W - total_w) // 2

def arch_box(x, title, sub_lines, tags, bg, br, tag_bg, tag_fg):
    add_rect(s, x, box_y, box_w, box_h, fill=bg, line=br,
             line_w=Pt(1.5), radius=0.06)
    add_text(s, x + px(20), box_y + px(44), box_w - px(40), px(44),
             title, size=fpt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + px(20), box_y + px(110), box_w - px(40), px(80),
             "\n".join(sub_lines), size=fpt(26),
             color=RGBColor(0x9a, 0xa3, 0xb0),
             align=PP_ALIGN.CENTER, line_spacing=1.4)
    tag_y = box_y + px(230); col_x = x + px(24)
    max_x = x + box_w - px(24)
    cur_x = col_x; row = 0
    for t in tags:
        tag_w = px(len(t) * 11 + 40)
        if cur_x + tag_w > max_x:
            row += 1; cur_x = col_x
        add_rect(s, cur_x, tag_y + row * px(50), tag_w, px(42),
                 fill=tag_bg, radius=0.5)
        add_text(s, cur_x, tag_y + row * px(50) + px(6), tag_w, px(32),
                 t, size=fpt(22), color=tag_fg, align=PP_ALIGN.CENTER)
        cur_x += tag_w + px(10)

arch_box(start_x, "React SPA",
         ["User Interface", "localhost:3000"],
         ["React 19", "Router", "recharts", "MetaMask", "jsPDF"],
         AB_REACT_BG, AB_REACT_BR, TAG_REACT_BG, TAG_REACT_FG)
arch_box(start_x + box_w + gap_x, "Node / Express API",
         ["Business Logic + MongoDB", "localhost:8000"],
         ["Express 4", "Mongoose", "JWT/bcrypt", "web3.js", "node-cron"],
         AB_NODE_BG, AB_NODE_BR, TAG_NODE_BG, TAG_NODE_FG)
arch_box(start_x + 2 * (box_w + gap_x), "Ganache Chain",
         ["Immutable Audit Layer", "localhost:7545"],
         ["Solidity 0.8.19", "Truffle v5", "Chain ID 5777"],
         AB_CHAIN_BG, AB_CHAIN_BR, TAG_CHAIN_BG, TAG_CHAIN_FG)

# Arrows + labels between boxes
def arrow_label(x_center, label_top, label_bot):
    add_text(s, x_center - px(110), box_y + box_h // 2 - px(70),
             px(220), px(36), label_top, size=fpt(22), bold=True,
             color=BLUE_LT, align=PP_ALIGN.CENTER)
    add_text(s, x_center - px(110), box_y + box_h // 2 + px(40),
             px(220), px(36), label_bot, size=fpt(20),
             color=RGBColor(0x9a, 0xa3, 0xb0), italic=True,
             align=PP_ALIGN.CENTER)
    line_y = box_y + box_h // 2
    add_arrow(s, x_center - px(28), line_y, x_center + px(28), line_y,
              color=BLUE_LT, weight=Pt(2.5))

arrow_label(start_x + box_w + gap_x // 2, "REST / JWT", "JSON + Bearer token")
arrow_label(start_x + 2 * box_w + gap_x + gap_x // 2, "web3.js / RPC",
            "signed tx + receipt")

add_text(s, px(120), px(950), px(1680), px(50),
         "Every booking, payment and prescription is mirrored on-chain — the tx hash is stored on the MongoDB document for independent verification.",
         size=fpt(24), color=RGBColor(0x8a, 0x94, 0xa2),
         align=PP_ALIGN.CENTER, italic=True)
set_notes(s, NOTES[2])

# --------------------------------------------------------------------------
# SLIDE 4 — Technology Stack
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "03 — TECHNOLOGIES", BLUE)
slide_title(s, "Technology Stack")

col_w = px(520); col_h = px(700); col_gap = px(40); col_x0 = px(120)

def stack_col(x, title, layer, items, bg, accent):
    add_rect(s, x, px(300), col_w, px(8), fill=accent)
    add_rect(s, x, px(308), col_w, col_h, fill=bg, radius=0.03)
    add_text(s, x + px(40), px(330), col_w - px(80), px(40),
             title, size=fpt(28), bold=True, color=accent, letter_spacing=200)
    # Layer badge
    layer_w = px(len(layer) * 11 + 40)
    add_rect(s, x + px(40), px(380), layer_w, px(36),
             fill=WHITE, line=accent, line_w=Pt(1), radius=0.5)
    add_text(s, x + px(40), px(384), layer_w, px(30),
             layer, size=fpt(20), bold=True, color=accent, align=PP_ALIGN.CENTER)
    y = px(440)
    for name, desc in items:
        add_text(s, x + px(40), y, col_w - px(80), px(42),
                 name, size=fpt(30), bold=True, color=TEXT)
        add_text(s, x + px(40), y + px(46), col_w - px(80), px(36),
                 desc, size=fpt(24), color=MUTED)
        y += px(116)

stack_col(col_x0, "FRONTEND", "PRESENTATION LAYER", [
    ("React 19 + Router v6", "SPA with role-based route protection"),
    ("recharts", "Admin analytics: revenue, diagnoses, on-chain coverage"),
    ("jsPDF", "Exportable bill & prescription PDFs"),
    ("web3.js + MetaMask", "window.ethereum — on-chain bill payment"),
], BG_BLUE_CARD, BLUE)

stack_col(col_x0 + col_w + col_gap, "BACKEND", "APPLICATION · PERSISTENCE", [
    ("Node.js + Express 4", "REST API, 12 route groups, Swagger UI"),
    ("MongoDB + Mongoose", "10 schemas + blockchainStatus audit fields"),
    ("JWT + bcryptjs", "Auth + doctor approval workflow"),
    ("web3.js + node-cron", "Contract calls & scheduled slot cleanup"),
], BG_TEAL_CARD, TEAL)

stack_col(col_x0 + 2 * (col_w + col_gap), "BLOCKCHAIN", "CONSENSUS LAYER", [
    ("Solidity 0.8.19", "HospitalRecords smart contract"),
    ("Truffle v5", "Compilation, migration, testing framework"),
    ("Ganache", "Local Ethereum chain — Chain ID 5777, :7545"),
    ("MetaMask", "Injected wallet for patient on-chain payments"),
], BG_AMBER_CARD, AMBER_DK)
set_notes(s, NOTES[3])

# --------------------------------------------------------------------------
# SLIDE 5 — Smart Contract Design (structs + code snippet)
# --------------------------------------------------------------------------
s = add_slide(NAVY)
slide_label(s, "04 — BLOCKCHAIN", BLUE_LT)
slide_title(s, "Smart Contract Design", color=WHITE)

# Left panel — structs & functions (compact)
lp_x = px(120); lp_y = px(300); lp_w = px(800); lp_h = px(670)
add_rect(s, lp_x, lp_y, lp_w, lp_h,
         fill=RGBColor(0x17, 0x24, 0x38),
         line=RGBColor(0x2a, 0x3a, 0x52), line_w=Pt(1), radius=0.04)
add_text(s, lp_x + px(40), lp_y + px(30), lp_w - px(80), px(38),
         "DATA STRUCTURES & CORE FUNCTIONS", size=fpt(24), bold=True,
         color=BLUE_LT, letter_spacing=180)

structs = [
    ("PatientRecord", "dataHash + timestamp + exists"),
    ("AppointmentRecord", "patientId · doctorId · status"),
    ("BillingRecord", "amount (wei) · status · paidBy"),
    ("PrescriptionRecord", "medicineHash (SHA-256)"),
]
cy = lp_y + px(90)
for name, desc in structs:
    add_text(s, lp_x + px(40), cy, px(360), px(34),
             name, size=fpt(24), bold=True, color=TEAL_LT, font="DM Mono")
    add_text(s, lp_x + px(40) + px(380), cy + px(4), px(360), px(34),
             desc, size=fpt(22), color=RGBColor(0x9a, 0xa3, 0xb0))
    cy += px(52)

# separator
add_rect(s, lp_x + px(40), cy + px(10), lp_w - px(80), px(1),
         fill=RGBColor(0x2a, 0x3a, 0x52))
cy += px(30)

add_text(s, lp_x + px(40), cy, lp_w - px(80), px(34),
         "KEY FUNCTIONS", size=fpt(22), bold=True,
         color=BLUE_LT, letter_spacing=180)
cy += px(50)

fns = [
    ("logAppointment(id, pat, dr, status)", "onlyOwner", BLUE_LT),
    ("logBilling(id, pat, amount, status)", "onlyOwner", BLUE_LT),
    ("logPrescription(id, pat, dr, hash)", "onlyOwner", BLUE_LT),
    ("payBillOnChain(billId, patientId)", "payable", AMBER),
    ("get*(id) — direct chain read", "public view", TEAL_LT),
]
for sig, mod, color in fns:
    add_text(s, lp_x + px(40), cy, px(560), px(32),
             sig, size=fpt(21), color=AMBER, font="DM Mono")
    add_rect(s, lp_x + px(620), cy + px(4), px(140), px(28),
             fill=RGBColor(0x1a, 0x2e, 0x48), radius=0.4)
    add_text(s, lp_x + px(620), cy + px(6), px(140), px(26),
             mod, size=fpt(18), color=color, font="DM Mono",
             align=PP_ALIGN.CENTER, bold=True)
    cy += px(48)

# Right panel — Solidity code snippet (the real payBillOnChain)
rp_x = px(940); rp_y = px(300); rp_w = px(860); rp_h = px(670)
add_rect(s, rp_x, rp_y, rp_w, rp_h,
         fill=CODE_BG, line=RGBColor(0x2a, 0x3a, 0x52), line_w=Pt(1),
         radius=0.04)
# File label bar
add_rect(s, rp_x, rp_y, rp_w, px(50),
         fill=RGBColor(0x15, 0x1f, 0x33), radius=0.08)
add_text(s, rp_x + px(24), rp_y + px(14), rp_w - px(48), px(28),
         "HospitalRecords.sol · payBillOnChain", size=fpt(22),
         color=BLUE_LT, font="DM Mono", bold=True)

code_lines = [
    ("function payBillOnChain(", CODE_FG),
    ("    string memory _billId,", CODE_FG),
    ("    string memory _patientId", CODE_FG),
    (") public payable {", CODE_FG),
    ("    require(msg.value > 0, \"Must send ETH\");", RGBColor(0xde, 0xbb, 0x6e)),
    ("    billingRecords[_billId].status = \"paid\";", CODE_FG),
    ("    billingRecords[_billId].paidBy = msg.sender;", CODE_FG),
    ("    payable(owner).transfer(msg.value);", RGBColor(0xb3, 0xe0, 0xa8)),
    ("    emit PaymentReceived(", CODE_FG),
    ("        _billId, msg.sender, msg.value, block.timestamp", CODE_FG),
    ("    );", CODE_FG),
    ("}", CODE_FG),
]
code_y = rp_y + px(80)
for line, color in code_lines:
    add_text(s, rp_x + px(24), code_y, rp_w - px(48), px(32),
             line, size=fpt(22), color=color, font="DM Mono",
             line_spacing=1.25)
    code_y += px(44)

add_text(s, rp_x + px(24), rp_y + rp_h - px(80), rp_w - px(48), px(60),
         "Patient → MetaMask → contract. No backend custody.\n"
         "msg.value is forwarded to owner; status flips inside the EVM.",
         size=fpt(20), color=TEAL_LT, italic=True, line_spacing=1.4)
set_notes(s, NOTES[4])

# --------------------------------------------------------------------------
# SLIDE 6 — Blockchain Interaction Flow (swimlane)
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "05 — INTERACTIONS", BLUE)
slide_title(s, "Blockchain Interaction Flow")
add_text(s, px(120), px(230), px(1680), px(40),
         "Example: Patient books an appointment — on-chain trail across five lanes",
         size=fpt(24), color=MUTED, italic=True)

# Swimlane layout: 5 vertical lanes
lane_y = px(300); lane_h = px(620)
lanes = [
    ("Patient",  RGBColor(0xd9, 0xa0, 0x4a)),
    ("React UI", BLUE),
    ("Express",  TEAL),
    ("MongoDB",  RGBColor(0x5a, 0xaf, 0x52)),
    ("Ganache",  RGBColor(0x14, 0x97, 0xb5)),
]
lane_x0 = px(120); lane_w_total = px(1680)
lane_w = lane_w_total // len(lanes)

for i, (label, color) in enumerate(lanes):
    x = lane_x0 + i * lane_w
    # header
    add_rect(s, x, lane_y, lane_w, px(60), fill=color, radius=0.05)
    add_text(s, x, lane_y + px(14), lane_w, px(34),
             label.upper(), size=fpt(24), bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, letter_spacing=200)
    # lane background (alternating)
    bg = ROW_ALT if i % 2 else WHITE
    add_rect(s, x, lane_y + px(60), lane_w, lane_h - px(60),
             fill=bg, line=BORDER, line_w=Pt(0.5))

# Steps positioned in lanes (lane_index, y_offset_from_lane_top, step_num, label)
steps = [
    (0, px(110), 1, "Fills booking\nform",
        "POST /api/\nappointments"),
    (1, px(110), 2, "JWT-auth'd\nrequest",
        "Bearer token\nverified"),
    (2, px(220), 3, "Creates\nAppointment",
        "Mongoose\nschema write"),
    (3, px(220), 4, "Document\npersisted",
        "ObjectId\nreturned"),
    (2, px(330), 5, "logAppointment()",
        "web3.send\nsigned tx"),
    (4, px(330), 6, "Tx mined,\nhash returned",
        "block.timestamp\nrecorded"),
    (3, px(440), 7, "txHash +\nstatus='logged'",
        "back on\nMongo doc"),
    (1, px(440), 8, "Notification\npolled (30s)",
        "bell icon\nflashes"),
]
for lane_i, y_off, num, title, detail in steps:
    x = lane_x0 + lane_i * lane_w + px(20)
    y = lane_y + y_off
    # circle
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + lane_w // 2 - px(80), y,
                           px(44), px(44))
    c.shadow.inherit = False; c.line.fill.background()
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    tf = c.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    rr = p.add_run(); rr.text = str(num)
    rr.font.size = fpt(22); rr.font.bold = True
    rr.font.color.rgb = WHITE; rr.font.name = 'DM Sans'
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # title + detail
    add_text(s, x + lane_w // 2 - px(30), y + px(2), lane_w - px(40), px(54),
             title, size=fpt(20), bold=True, color=TEXT, line_spacing=1.15)
    add_text(s, x + lane_w // 2 - px(30), y + px(54), lane_w - px(40), px(50),
             detail, size=fpt(18), color=MUTED, line_spacing=1.25)

# Arrows between consecutive steps (simple — small arrow at mid of lane boundary)
arrow_pairs = [(0, 1), (1, 2), (2, 3), (3, 4), (4, 5), (5, 6), (6, 7)]
for a, b in arrow_pairs:
    la, ya, _, _, _ = steps[a]
    lb, yb, _, _, _ = steps[b]
    x1 = lane_x0 + la * lane_w + lane_w // 2 + px(20)
    x2 = lane_x0 + lb * lane_w + lane_w // 2 - px(60)
    y1 = lane_y + ya + px(22); y2 = lane_y + yb + px(22)
    add_arrow(s, x1, y1, x2, y2, color=MUTED, weight=Pt(1.5))

set_notes(s, NOTES[5])

# --------------------------------------------------------------------------
# SLIDE 7 — Off-chain vs On-chain Payment (NEW)
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "06 — PAYMENT PATHS", BLUE)
slide_title(s, "Off-Chain vs On-Chain Payment")
add_text(s, px(120), px(230), px(1680), px(40),
         "Two payment flows coexist — the right path is what makes this a DApp",
         size=fpt(24), color=MUTED, italic=True)

def payment_path(x, title, badge, badge_bg, badge_fg, steps, accent):
    panel_w = px(820); panel_h = px(680); panel_y = px(300)
    add_rect(s, x, panel_y, panel_w, panel_h,
             fill=ROW_ALT, line=accent, line_w=Pt(2), radius=0.03)
    # header
    add_rect(s, x, panel_y, panel_w, px(100),
             fill=accent, radius=0.08)
    add_text(s, x + px(40), panel_y + px(20), panel_w - px(80), px(40),
             title, size=fpt(34), bold=True, color=WHITE)
    # badge
    badge_w = px(len(badge) * 12 + 40)
    add_rect(s, x + px(40), panel_y + px(62), badge_w, px(30),
             fill=badge_bg, radius=0.5)
    add_text(s, x + px(40), panel_y + px(66), badge_w, px(26),
             badge, size=fpt(20), bold=True, color=badge_fg,
             align=PP_ALIGN.CENTER, letter_spacing=200)

    sy = panel_y + px(130)
    for i, (step, note) in enumerate(steps):
        # step number circle
        cx = x + px(40); cy = sy + px(6)
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, px(40), px(40))
        c.shadow.inherit = False; c.line.fill.background()
        c.fill.solid(); c.fill.fore_color.rgb = accent
        tf = c.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        rr = p.add_run(); rr.text = str(i + 1)
        rr.font.size = fpt(22); rr.font.bold = True
        rr.font.color.rgb = WHITE; rr.font.name = 'DM Sans'
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        add_text(s, x + px(100), sy, panel_w - px(140), px(36),
                 step, size=fpt(26), bold=True, color=TEXT)
        add_text(s, x + px(100), sy + px(36), panel_w - px(140), px(32),
                 note, size=fpt(22), color=MUTED, line_spacing=1.3)
        sy += px(82)

payment_path(px(120), "Backend-Mediated", "TRUSTS BACKEND",
             RGBColor(0xf3, 0xe3, 0xc2), AMBER_DK, [
    ("Patient clicks 'Pay'",
     "REST call to /api/bills/:id/pay"),
    ("Express marks bill paid in MongoDB",
     "bill.status = 'paid'"),
    ("Backend calls logBilling()",
     "web3.send from DEPLOYER_ACCOUNT"),
    ("Tx hash stored on Mongo doc",
     "blockchainStatus = 'logged'"),
    ("Audit middleware writes AuditLog",
     "who/what/when — off-chain"),
], AMBER_DK)

payment_path(px(980), "MetaMask On-Chain", "TRUSTLESS · DApp",
             RGBColor(0xc3, 0xed, 0xe1), TEAL, [
    ("Patient clicks 'Pay with MetaMask'",
     "window.ethereum.request"),
    ("MetaMask prompts to sign tx",
     "gas estimate shown, patient approves"),
    ("payBillOnChain() executes in EVM",
     "ETH flows patient → contract → owner"),
    ("Contract flips status, emits event",
     "PaymentReceived + BillingLogged"),
    ("Backend never holds patient funds",
     "attack surface removed"),
], TEAL)

set_notes(s, NOTES[6])

# --------------------------------------------------------------------------
# SLIDE 8 — Role-Based Features Table
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "07 — FEATURES", BLUE)
slide_title(s, "Role-Based Features")

table_x = px(120); table_y = px(270)
table_w = px(1680); table_h = px(720)
header_h = px(80)
col_widths = [px(380), px(310), px(330), px(330), px(330)]

add_rect(s, table_x, table_y, table_w, table_h, line=BORDER,
         line_w=Pt(1.5), radius=0.015)
add_rect(s, table_x, table_y, table_w, header_h, fill=NAVY, radius=0.015)

col_labels = [("", "", WHITE), ("ADMIN", "★", BLUE_LT),
              ("DOCTOR", "⚕", TEAL_LT),
              ("PATIENT", "♥", AMBER),
              ("RECEPTIONIST", "◈", RGBColor(0x9a, 0xa3, 0xb0))]
cx = table_x
for (lbl, icon, col), w in zip(col_labels, col_widths):
    if lbl:
        add_text(s, cx, table_y + px(14), w, px(34),
                 icon, size=fpt(28), color=col, align=PP_ALIGN.CENTER)
        add_text(s, cx, table_y + px(46), w, px(34),
                 lbl, size=fpt(22), bold=True, color=col,
                 align=PP_ALIGN.CENTER, letter_spacing=200)
    cx += w

rows_data = [
    ("Analytics Dashboard",    ("✓", 'yes'), ("—", 'no'), ("—", 'no'), ("—", 'no')),
    ("Blockchain Verification",("✓", 'yes'), ("—", 'no'), ("—", 'no'), ("—", 'no')),
    ("Manage Appointments",    ("✓ all", 'yes'), ("✓ own", 'yes'),
                               ("Book/Cancel", 'part'), ("—", 'no')),
    ("Prescriptions",          ("✓ view all", 'yes'), ("✓ create", 'yes'),
                               ("View own", 'part'), ("—", 'no')),
    ("Billing & Payment",      ("✓ view all", 'yes'), ("View own", 'part'),
                               ("✓ pay", 'yes'), ("—", 'no')),
    ("MetaMask On-Chain Pay",  ("—", 'no'), ("—", 'no'), ("✓", 'yes'), ("—", 'no')),
    ("Beds & Departments",     ("✓ manage", 'yes'), ("—", 'no'),
                               ("—", 'no'), ("✓ assign", 'yes')),
    ("Audit Log",              ("✓", 'yes'), ("—", 'no'), ("—", 'no'), ("—", 'no')),
    ("PDF Export",             ("✓", 'yes'), ("✓", 'yes'),
                               ("✓", 'yes'), ("—", 'no')),
]
row_h = (table_h - header_h) / len(rows_data)
for i, row in enumerate(rows_data):
    ry = table_y + header_h + i * row_h
    if i % 2 == 1:
        add_rect(s, table_x, ry, table_w, row_h, fill=ROW_ALT)
    if i > 0:
        sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, table_x, ry,
                                 table_w, px(1))
        sep.shadow.inherit = False; sep.line.fill.background()
        sep.fill.solid(); sep.fill.fore_color.rgb = BORDER
    cx = table_x
    feature = row[0]
    add_text(s, cx + px(20), ry, col_widths[0] - px(40), row_h,
             feature, size=fpt(24), bold=True, color=TEXT,
             anchor=MSO_ANCHOR.MIDDLE)
    cx += col_widths[0]
    for j, (val, kind) in enumerate(row[1:]):
        col = TEAL if kind == 'yes' else \
              (AMBER_DK if kind == 'part' else RGBColor(0xb8, 0xbe, 0xca))
        add_text(s, cx, ry, col_widths[j + 1], row_h,
                 val, size=fpt(24),
                 bold=(kind == 'yes'), color=col,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        cx += col_widths[j + 1]
set_notes(s, NOTES[7])

# --------------------------------------------------------------------------
# SLIDE 9 — Live Screenshots (NEW)
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "08 — THE RUNNING APP", BLUE)
slide_title(s, "Live Screenshots")

screenshots = [
    ("shot-01-admin-dashboard.png",        "Admin — overview + recharts analytics"),
    ("shot-02-appointment-booking.png",    "Patient — booking flow"),
    ("shot-03-metamask-payment.png",       "MetaMask — payBillOnChain signature"),
    ("shot-04-blockchain-verification.png","Admin — on-chain verification"),
    ("shot-05-audit-log.png",              "Admin — audit trail"),
    ("shot-06-prescription-pdf.png",       "jsPDF — portable record"),
]

grid_x0 = px(120); grid_y0 = px(280)
cell_w = px(540); cell_h = px(340)
gap_x = px(20); gap_y = px(30)

for i, (fname, caption) in enumerate(screenshots):
    col = i % 3; row = i // 3
    x = grid_x0 + col * (cell_w + gap_x)
    y = grid_y0 + row * (cell_h + gap_y)
    # placeholder / image
    img_h = cell_h - px(54)
    add_rect(s, x, y, cell_w, img_h, fill=ROW_ALT,
             line=BORDER, line_w=Pt(1), radius=0.02)
    fpath = SCREENSHOT_DIR / fname
    if fpath.exists():
        try:
            s.shapes.add_picture(str(fpath), x + px(4), y + px(4),
                                 width=cell_w - px(8), height=img_h - px(8))
        except Exception as e:
            print(f"[warn] failed to embed {fname}: {e}")
    else:
        add_text(s, x, y + img_h // 2 - px(20), cell_w, px(40),
                 f"[ {fname} ]", size=fpt(22), color=MUTED,
                 align=PP_ALIGN.CENTER, italic=True)
    add_text(s, x, y + img_h + px(8), cell_w, px(36),
             caption, size=fpt(22), color=TEXT,
             align=PP_ALIGN.CENTER, bold=True)

set_notes(s, NOTES[8])

# --------------------------------------------------------------------------
# SLIDE 10 — Cross-Cutting Features
# --------------------------------------------------------------------------
s = add_slide(NAVY)
slide_label(s, "09 — CROSS-CUTTING", BLUE_LT)
slide_title(s, "Key Cross-Cutting Features", color=WHITE)

feats = [
    (BLUE_LT, "Tamper-Proof Audit Trail",
     "Every appointment, billing, and prescription is mirrored on-chain. The tx hash is persisted on the Mongo document; admins re-verify by reading directly from Ganache."),
    (TEAL_LT, "Audit Middleware",
     "An Express middleware intercepts every mutating HTTP request and writes a timestamped AuditLog entry — who, what, when — independent of the blockchain trail."),
    (AMBER, "MetaMask Wallet Integration",
     "The MetaMask bar detects window.ethereum, shows account + ETH balance. Patients call payBillOnChain directly — backend never holds their funds."),
    (DOT_PURPLE, "Real-Time Notifications",
     "In-app bell polls the API every 30 s. Events fire on booking confirmation, payment success, and prescription issuance — no page refresh required."),
    (DOT_RED, "PDF Exports via jsPDF",
     "Bills and prescriptions exported from patient and doctor dashboards, with the tx hash embedded in the footer so paper records remain verifiable."),
    (DOT_GREEN, "JWT Auth + Doctor Approval",
     "Doctor registrations enter a 'pending' state requiring admin approval. bcryptjs + JWT secure every endpoint; role middleware enforces access boundaries."),
]
fg_x0 = px(120); fg_y0 = px(290)
fc_w = px(530); fc_h = px(330); fc_gap = px(25)
for i, (dot, title, body) in enumerate(feats):
    col = i % 3; row = i // 3
    x = fg_x0 + col * (fc_w + fc_gap); y = fg_y0 + row * (fc_h + fc_gap)
    add_rect(s, x, y, fc_w, fc_h,
             fill=RGBColor(0x1a, 0x29, 0x3e),
             line=RGBColor(0x2a, 0x3a, 0x52), line_w=Pt(1), radius=0.05)
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, x + px(36), y + px(36),
                           px(18), px(18))
    d.shadow.inherit = False; d.line.fill.background()
    d.fill.solid(); d.fill.fore_color.rgb = dot
    add_text(s, x + px(36), y + px(72), fc_w - px(72), px(48),
             title, size=fpt(30), bold=True, color=WHITE)
    add_text(s, x + px(36), y + px(130), fc_w - px(72), fc_h - px(150),
             body, size=fpt(22), color=RGBColor(0x98, 0xa1, 0xae),
             line_spacing=1.45)
set_notes(s, NOTES[9])

# --------------------------------------------------------------------------
# SLIDE 11 — Design Rationale & Critical Evaluation (NEW)
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "10 — CRITICAL ANALYSIS", BLUE)
slide_title(s, "Design Rationale & Evaluation")

col_w = px(540); col_h = px(680); col_gap = px(24); col_x0 = px(120)

rationale_cols = [
    ("WHY THESE CHOICES", BLUE, BG_BLUE_CARD, [
        ("Ganache over Sepolia",
         "Zero-cost, instant feedback loop during development — no faucet or block-time delay."),
        ("Mongoose over raw driver",
         "Schema validation + middleware hooks enabled audit logging with zero controller noise."),
        ("Hybrid, not fully on-chain",
         "Storing full records on-chain is prohibitively expensive; we anchor hashes + status only."),
        ("JWT over sessions",
         "Stateless auth scales better and matches the SPA model — no server-side session store."),
    ]),
    ("TRADE-OFFS & LIMITATIONS", AMBER_DK, BG_AMBER_CARD, [
        ("onlyOwner centralisation",
         "A single deployer account signs every write — reintroduces trust the chain was meant to remove."),
        ("Ganache is local",
         "Tutors cannot inspect the chain on a public explorer; demo-only, not production-ready."),
        ("No event indexer",
         "Historical queries iterate mappings — won't scale past thousands of records."),
        ("Silent errors fixed, but incomplete",
         "We now persist 'failed' status; a retry queue would be the next step."),
    ]),
    ("FUTURE WORK", TEAL, BG_TEAL_CARD, [
        ("Migrate to Sepolia testnet",
         "Public verifiability + Etherscan integration for tx audit."),
        ("TheGraph event indexer",
         "Scalable historical queries and real-time dashboards."),
        ("Multi-sig deployer",
         "Decentralise write authority across hospital stakeholders."),
        ("Expand Jest + add E2E",
         "Coverage beyond the 3 coursework suites — Cypress for the React flows."),
    ]),
]

for i, (title, accent, bg, items) in enumerate(rationale_cols):
    x = col_x0 + i * (col_w + col_gap)
    add_rect(s, x, px(280), col_w, px(8), fill=accent)
    add_rect(s, x, px(288), col_w, col_h, fill=bg, radius=0.03)
    add_text(s, x + px(30), px(310), col_w - px(60), px(40),
             title, size=fpt(24), bold=True, color=accent, letter_spacing=200)
    cy = px(370)
    for name, desc in items:
        add_text(s, x + px(30), cy, col_w - px(60), px(40),
                 name, size=fpt(24), bold=True, color=TEXT)
        add_text(s, x + px(30), cy + px(42), col_w - px(60), px(110),
                 desc, size=fpt(20), color=MUTED, line_spacing=1.4)
        cy += px(150)

set_notes(s, NOTES[10])

# --------------------------------------------------------------------------
# SLIDE 12 — Baseline vs Extension (master vs main)
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "11 — SCOPE OF CONTRIBUTION", BLUE)
slide_title(s, "Baseline vs Extension")

# Two-column comparison
col_w = px(820); col_h = px(560); col_gap = px(40); col_x0 = px(120)
col_y = px(280)

baseline_items = [
    ("Node.js + Express API",        "Auth, appointments, bills — REST only."),
    ("MongoDB + Mongoose",            "Schemas for User, Appointment, Bill, Prescription."),
    ("JWT + bcrypt auth",             "Role-agnostic, no doctor-approval flow."),
    ("Swagger docs + seed script",    "Basic Postman-style API surface."),
    ("No frontend",                   "Backend-only — no UI, no dashboards."),
    ("No blockchain",                 "All trust on the API layer."),
]

extension_items = [
    ("React 19 SPA",                  "Role-based dashboards for admin / doctor / patient / receptionist."),
    ("HospitalRecords.sol",           "Solidity 0.8.19 contract — 4 structs, payBillOnChain payable."),
    ("web3.js + Ganache pipeline",    "Every mutating write mirrored on-chain; tx hash persisted to Mongo."),
    ("MetaMask integration",          "Fully decentralised patient payments via window.ethereum."),
    ("recharts analytics + jsPDF",    "Admin charts, revenue trends, exportable bill/rx PDFs."),
    ("Audit log, notifications, tests","Mutation middleware, bell polling, Jest suites, RBAC fixes."),
]

def compare_col(x, title, accent, bg, items):
    # accent bar
    add_rect(s, x, col_y, col_w, px(8), fill=accent)
    # panel
    add_rect(s, x, col_y + px(8), col_w, col_h, fill=bg, radius=0.03)
    # title
    add_text(s, x + px(30), col_y + px(30), col_w - px(60), px(44),
             title, size=fpt(26), bold=True, color=accent, letter_spacing=200)
    cy = col_y + px(90)
    for name, desc in items:
        add_text(s, x + px(30), cy, col_w - px(60), px(32),
                 name, size=fpt(22), bold=True, color=TEXT)
        add_text(s, x + px(30), cy + px(34), col_w - px(60), px(50),
                 desc, size=fpt(18), color=MUTED, line_spacing=1.3)
        cy += px(80)

compare_col(col_x0, "MASTER BRANCH — BASELINE", MUTED, BG_AMBER_CARD, baseline_items)
compare_col(col_x0 + col_w + col_gap, "MAIN BRANCH — SUBMISSION", BLUE, BG_BLUE_CARD, extension_items)

# Footer — repo URLs
footer_y = col_y + col_h + px(40)
add_rect(s, px(120), footer_y, SLIDE_W - px(240), px(90),
         fill=RGBColor(0xf2, 0xf5, 0xfa), radius=0.04)
add_text(s, px(140), footer_y + px(12), SLIDE_W - px(280), px(28),
         "REPO LINKS", size=fpt(18), bold=True, color=BLUE,
         letter_spacing=200)
add_text(s, px(140), footer_y + px(40), SLIDE_W - px(280), px(44),
         "Submission (main):  github.com/rahmane001/hospital-system/tree/main\n"
         "Baseline (master): github.com/rahmane001/hospital-system/tree/master   ·   "
         "Diff: github.com/rahmane001/hospital-system/compare/master...main",
         size=fpt(16), color=TEXT, font="DM Mono", line_spacing=1.4)

set_notes(s, NOTES[11])

# --------------------------------------------------------------------------
# SLIDE 13 — Video Demonstration
# --------------------------------------------------------------------------
s = add_slide(NAVY)
add_text(s, 0, px(260), SLIDE_W, px(34),
         "12 — DEMONSTRATION", size=fpt(24), bold=True, color=BLUE_LT,
         align=PP_ALIGN.CENTER, letter_spacing=240)
add_text(s, 0, px(330), SLIDE_W, px(120),
         "Video Demonstration", size=fpt(80), bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
box_w = px(1200); box_h = px(220)
bx = (SLIDE_W - box_w) // 2; by = px(530)
add_rect(s, bx, by, box_w, box_h,
         fill=RGBColor(0x18, 0x2a, 0x42), line=RGBColor(0x24, 0x4a, 0x70),
         line_w=Pt(2), radius=0.08)
add_text(s, bx, by + px(40), box_w, px(36),
         "MS STREAM RECORDING", size=fpt(24),
         color=RGBColor(0x8a, 0x94, 0xa2), align=PP_ALIGN.CENTER,
         letter_spacing=200, bold=True)
add_text(s, bx + px(40), by + px(100), box_w - px(80), px(80),
         "[ Insert MS Stream link here ]",
         size=fpt(32), color=BLUE_LT, font="DM Mono",
         align=PP_ALIGN.CENTER)
add_text(s, px(200), by + box_h + px(60), SLIDE_W - px(400), px(130),
         "Access has been granted to all lab tutors.\n"
         "The recording covers login as each role · appointment booking · "
         "on-chain transaction logging · MetaMask bill payment · "
         "blockchain verification.",
         size=fpt(24), color=RGBColor(0x8a, 0x94, 0xa2),
         align=PP_ALIGN.CENTER, line_spacing=1.6)
set_notes(s, NOTES[12])

# --------------------------------------------------------------------------
# SLIDE 14 — References (Harvard, updated to 2026)
# --------------------------------------------------------------------------
s = add_slide(WHITE)
slide_label(s, "13 — REFERENCES", BLUE)
slide_title(s, "References")

refs = [
    ("Antonopoulos, A.M. and Wood, G. (2018) ",
     "Mastering Ethereum: Building Smart Contracts and DApps.",
     " Sebastopol, CA: O'Reilly Media."),
    ("Tanenbaum, A.S. and Van Steen, M. (2017) ",
     "Distributed Systems: Principles and Paradigms.",
     " 3rd edn. CreateSpace."),
    ("Wood, G. (2014) ",
     "Ethereum: A Secure Decentralised Generalised Transaction Ledger (Yellow Paper).",
     " Available at: https://ethereum.github.io/yellowpaper/paper.pdf [Accessed: April 2026]."),
    ("Ethereum Foundation (2024) ",
     "Solidity Documentation v0.8.19.",
     " Available at: https://docs.soliditylang.org/en/v0.8.19/ [Accessed: April 2026]."),
    ("Truffle Suite (2024) ",
     "Truffle Documentation.",
     " Available at: https://trufflesuite.com/docs/ [Accessed: April 2026]."),
    ("web3.js Contributors (2024) ",
     "web3.js Documentation.",
     " Available at: https://docs.web3js.org/ [Accessed: April 2026]."),
    ("Nakamoto, S. (2008) 'Bitcoin: A Peer-to-Peer Electronic Cash System.' ",
     "bitcoin.org.",
     " Available at: https://bitcoin.org/bitcoin.pdf [Accessed: April 2026]."),
    ("MongoDB, Inc. (2024) ",
     "Mongoose ODM Documentation.",
     " Available at: https://mongoosejs.com/docs/ [Accessed: April 2026]."),
]

ref_y = px(270); ref_x = px(150); ref_w = px(1620)
ref_h = px(82); ref_gap = px(6)
for i, (pre, em, post) in enumerate(refs):
    y = ref_y + i * (ref_h + ref_gap)
    add_rect(s, ref_x, y + px(10), px(4), ref_h - px(20), fill=BLUE)
    add_text(s, ref_x - px(40), y + px(20), px(36), px(40),
             f"[{i+1}]", size=fpt(20), color=MUTED, align=PP_ALIGN.RIGHT)
    add_multiline(s, ref_x + px(28), y + px(6), ref_w - px(60), ref_h,
        [[
            {'text': pre, 'size': fpt(22), 'color': TEXT},
            {'text': em, 'size': fpt(22), 'color': TEXT, 'italic': True},
            {'text': post, 'size': fpt(22), 'color': TEXT},
        ]],
        line_spacing=1.4, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, NOTES[13])

# --------------------------------------------------------------------------
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
missing = [f for f, _ in screenshots if not (SCREENSHOT_DIR / f).exists()]
if missing:
    print(f"[warn] {len(missing)} screenshot(s) missing from "
          f"{SCREENSHOT_DIR}: {missing}")
